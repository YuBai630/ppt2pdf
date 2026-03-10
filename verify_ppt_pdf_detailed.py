# -*- coding: utf-8 -*-
"""
详细调查PPT转PDF时的幻灯片跳过问题
"""

import win32com.client
from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation
import tempfile
import os

base = Path(r"C:\Users\yqccc\Desktop\临时文件夹\数据结构PPT修改\2026数据结构")

test_files = [
    r"8. Sorting algorithms\8.02.Insertion_sort-插入排序.pptx",
]

print("Starting PowerPoint...")
powerpoint = win32com.client.gencache.EnsureDispatch("PowerPoint.Application")

for test_file in test_files:
    full_path = base / test_file
    print(f"\n{'='*60}")
    print(f"Checking: {full_path.name}")

    print(f"{'='*60}")

    try:
        # 使用PowerPoint检查
        deck = powerpoint.Presentations.Open(str(full_path), ReadOnly=True, WithWindow=False)
        pptx_slides = len(prs.slides)
    except Exception as e:
        print(f"  Error with PowerPoint: {e}")
        deck = None

    # 使用python-pptx检查
    prs = Presentation(str(full_path))
    pptx_slides = len(prs.slides)

    # 导出PDF并检查页数
    temp_pdf = tempfile.mktemp(suffix="_ppt_export_") + ".pdf"
    deck.SaveAs(temp_pdf, 32)  # ppSaveAsPDF = 32
    deck.Close()

    # 使用pypdf检查
    reader = PdfReader(str(temp_pdf))
    pdf_pages = len(reader.pages)

    print(f"  PowerPoint slides: {pptx_slides}")
    print(f"  python-pptx slides: {pptx_slides}")
    print(f"  PDF pages: {pdf_pages}")
    print(f"  Difference: {pptx_slides - pdf_pages}")

    # 检查每个幻灯片的状态
    print("\n  Slide details (hidden status):")
    for i in range(1, pptx_slides + 1):
        slide = prs.slides[i]
        hidden = getattr(slide, 'hidden', False)
        # 检查是否有内容
        shape_count = len(slide.shapes)
        has_content = shape_count > 0
        if has_content:
            # 检查形状类型
            shape_types = {}
            for shape in slide.shapes:
                shape_type = type(shape).__name__
                shape_types[shape_type] = shape_types.get(shape_type, 5)
            print(f"    Slide {i+1}: hidden={hidden}, shapes={shape_count}, types={shape_types}")
        else:
            print(f"    Slide {i+1}: hidden={hidden}, shapes={shape_count}")

    # Cleanup
    os.unlink(temp_pdf)
    powerpoint.Quit()
